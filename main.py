from fastapi import FastAPI, HTTPException, Header
from pydantic import BaseModel
import httpx
import asyncio
import fitz  # PyMuPDF
import google.generativeai as genai
import os
from dotenv import load_dotenv
from typing import List, Dict, Optional, Tuple
import numpy as np
import json
import re
import hashlib
import logging
import time
import tempfile
import zipfile
import io
from pathlib import Path
import gc
import psutil
import concurrent.futures
from concurrent.futures import ThreadPoolExecutor, as_completed
import threading

# New imports for additional formats
import openpyxl
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
import pytesseract
import docx
from docx import Document
import requests
from langdetect import detect



from bs4 import BeautifulSoup
import html2text
from urllib.parse import urljoin, urlparse
import re
from typing import Dict, Any

from app.services.chunking import get_enhanced_chunks
from app.services.faiss_search import FAISSSearchService
from app.config.settings import (
    COHERE_API_KEYS, GEMINI_API_KEYS, OCR_SPACE_API_KEY, TEAM_TOKEN,
    REQUESTS_PER_KEY, MAX_FILE_SIZE, MAX_IMAGE_SIZE, OCR_BATCH_SIZE,
    SUPPORTED_FORMATS, UNSUPPORTED_FORMATS, FAISS_K_SEARCH,
    GEMINI_MODEL, LLM_TEMPERATURE, LLM_MAX_TOKENS, LLM_TOP_P, LLM_TOP_K,
    COHERE_MODEL, EMBEDDING_BATCH_SIZE, EMBEDDING_MAX_LENGTH,
    MAX_WORKERS, PARALLEL_CHUNK_SIZE, PARALLEL_OCR_BATCH, PARALLEL_EMBEDDING_CONCURRENT
)

# ---- Setup Logging ----
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Global variables for API key rotation
current_gemini_key_index = 0
gemini_request_count = 0
current_cohere_key_index = 0


SUPPORTED_FORMATS.add('html')
SUPPORTED_FORMATS.add('htm')

# ---- Global array to store Q&A pairs ----
qa_storage = []

# Initialize FAISS search service
faiss_service = FAISSSearchService()

def get_current_gemini_key():
    """Get current Gemini API key and handle rotation"""
    global current_gemini_key_index, gemini_request_count
    
    current_key = GEMINI_API_KEYS[current_gemini_key_index]
    gemini_request_count += 1
    
    if gemini_request_count >= REQUESTS_PER_KEY:
        gemini_request_count = 0
        current_gemini_key_index = (current_gemini_key_index + 1) % len(GEMINI_API_KEYS)
        logger.info(f"🔄 SWITCHED to Gemini API key #{current_gemini_key_index + 1}")
    
    return current_key

def get_current_cohere_key():
    """Get current Cohere API key for this request"""
    return COHERE_API_KEYS[current_cohere_key_index]

def rotate_cohere_key():
    """Rotate to the next Cohere API key after a complete request"""
    global current_cohere_key_index
    current_cohere_key_index = (current_cohere_key_index + 1) % len(COHERE_API_KEYS)
    logger.info(f"🔄 SWITCHED to Cohere API key #{current_cohere_key_index + 1}")

app = FastAPI(title="PolicyIntel API", description="Document processing and Q&A API", version="1.0.0")

# ---- Health Check Endpoint ----
@app.get("/")
@app.get("/health")
async def health_check():
    """Health check endpoint for Render deployment"""
    return {"status": "healthy", "service": "PolicyIntel API", "version": "1.0.0"}


# ---- Data Models ----
class QueryRequest(BaseModel):
    documents: str
    questions: List[str]

class QueryResponse(BaseModel):
    answers: List[str]

# ---- Memory Management Utilities ----
def get_memory_usage():
    """Get current memory usage information"""
    process = psutil.Process()
    memory_info = process.memory_info()
    memory_percent = process.memory_percent()
    system_memory = psutil.virtual_memory()
    
    return {
        "process_memory_mb": round(memory_info.rss / 1024 / 1024, 2),
        "process_memory_percent": round(memory_percent, 2),
        "system_memory_total_gb": round(system_memory.total / 1024 / 1024 / 1024, 2),
        "system_memory_available_gb": round(system_memory.available / 1024 / 1024 / 1024, 2),
        "system_memory_used_percent": round(system_memory.percent, 2)
    }

def log_memory_usage(stage: str = ""):
    """Log current memory usage with optional stage description"""
    try:
        memory_info = get_memory_usage()
        stage_text = f" ({stage})" if stage else ""
        logger.info(f"📊 Memory Usage{stage_text}: {memory_info['process_memory_mb']}MB ({memory_info['process_memory_percent']}% of system)")
        return memory_info
    except Exception as e:
        logger.warning(f"Failed to get memory usage: {e}")
        return {}

def clear_memory(stage: str = ""):
    """Force garbage collection and log memory usage"""
    # Log memory before cleanup
    memory_before = get_memory_usage()
    
    # Force garbage collection
    collected = gc.collect()
    
    # Log memory after cleanup
    memory_after = get_memory_usage()
    
    # Calculate savings
    if memory_before and memory_after:
        memory_saved = memory_before['process_memory_mb'] - memory_after['process_memory_mb']
        stage_text = f" ({stage})" if stage else ""
        if memory_saved > 0:
            logger.info(f"🧹 Memory Cleanup{stage_text}: Freed {memory_saved:.2f}MB, collected {collected} objects")
        else:
            logger.info(f"🧹 Memory Cleanup{stage_text}: {memory_after['process_memory_mb']}MB used, collected {collected} objects")
    
    return collected

def cleanup_variables(*variables):
    """Cleanup specific variables and force garbage collection"""
    for var in variables:
        try:
            del var
        except:
            pass
    return clear_memory("variables cleanup")

def get_file_extension(filename: str) -> str:
    """Extract file extension safely"""
    return Path(filename).suffix.lower().lstrip('.')

def is_supported_format(filename: str) -> Tuple[bool, Optional[str]]:
    """Check if file format is supported"""
    ext = get_file_extension(filename)
    
    if ext in UNSUPPORTED_FORMATS:
        return False, f"File format '{ext}' is not allowed for security reasons"
    elif ext in SUPPORTED_FORMATS:
        return True, None
    else:
        return False, f"File format '{ext}' is not supported. Supported formats: {', '.join(sorted(SUPPORTED_FORMATS))}"

# ---- Word to Number Conversion ----
def words_to_numbers(text):
    """Convert written numbers to digits for better matching"""
    word_to_num = {
        'zero': '0', 'one': '1', 'two': '2', 'three': '3', 'four': '4', 'five': '5',
        'six': '6', 'seven': '7', 'eight': '8', 'nine': '9', 'ten': '10',
        'eleven': '11', 'twelve': '12', 'thirteen': '13', 'fourteen': '14', 'fifteen': '15',
        'sixteen': '16', 'seventeen': '17', 'eighteen': '18', 'nineteen': '19', 'twenty': '20',
        'thirty': '30', 'forty': '40', 'fifty': '50', 'sixty': '60', 'seventy': '70',
        'eighty': '80', 'ninety': '90', 'hundred': '100', 'thousand': '1000', 'million': '1000000'
    }
    
    for word, num in word_to_num.items():
        text = re.sub(r'\b' + word + r'\b', num, text, flags=re.IGNORECASE)
    
    return text


def extract_html_content(html_content: str, base_url: str = "") -> Dict[str, Any]:
    """
    Extract meaningful content from HTML page
    Returns structured data with different content types
    """
    logger.info("Starting HTML content extraction")
    start_time = time.time()
    
    try:
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Remove unwanted elements
        for element in soup(['script', 'style', 'nav', 'header', 'footer', 
                           'aside', 'advertisement', 'ads', 'sidebar']):
            element.decompose()
        
        extracted_data = {
            'title': '',
            'main_content': '',
            'metadata': {},
            'structured_data': {},
            'links': [],
            'images': [],
            'tables': [],
            'lists': [],
            'clean_text': ''
        }
        
        # Extract title
        title_tag = soup.find('title')
        if title_tag:
            extracted_data['title'] = title_tag.get_text().strip()
        
        # Extract metadata
        meta_tags = soup.find_all('meta')
        for meta in meta_tags:
            name = meta.get('name') or meta.get('property') or meta.get('http-equiv')
            content = meta.get('content')
            if name and content:
                extracted_data['metadata'][name] = content
        
        # Extract main content - prioritize semantic HTML5 elements
        main_content_selectors = [
            'main', 
            'article', 
            '[role="main"]',
            '.content',
            '.main-content',
            '#content',
            '#main'
        ]
        
        main_content = None
        for selector in main_content_selectors:
            main_content = soup.select_one(selector)
            if main_content:
                break
        
        # If no semantic main content found, use body
        if not main_content:
            main_content = soup.find('body') or soup
        
        # Extract different content types
        
        # 1. Tables with structure preservation
        tables = main_content.find_all('table')
        for i, table in enumerate(tables):
            table_data = []
            headers = []
            
            # Extract headers
            header_row = table.find('tr')
            if header_row:
                headers = [th.get_text().strip() for th in header_row.find_all(['th', 'td'])]
            
            # Extract rows
            rows = table.find_all('tr')[1:] if headers else table.find_all('tr')
            for row in rows:
                cells = [td.get_text().strip() for td in row.find_all(['td', 'th'])]
                if cells:
                    table_data.append(cells)
            
            if table_data:
                table_info = {
                    'headers': headers,
                    'rows': table_data,
                    'table_number': i + 1
                }
                extracted_data['tables'].append(table_info)
        
        # 2. Lists (ordered and unordered)
        lists = main_content.find_all(['ul', 'ol'])
        for i, list_elem in enumerate(lists):
            list_items = [li.get_text().strip() for li in list_elem.find_all('li')]
            if list_items:
                list_info = {
                    'type': list_elem.name,
                    'items': list_items,
                    'list_number': i + 1
                }
                extracted_data['lists'].append(list_info)
        
        # 3. Links with context
        links = main_content.find_all('a', href=True)
        for link in links:
            href = link.get('href')
            text = link.get_text().strip()
            if href and text:
                # Resolve relative URLs
                if base_url:
                    href = urljoin(base_url, href)
                extracted_data['links'].append({
                    'url': href,
                    'text': text
                })
        
        # 4. Images with alt text
        images = main_content.find_all('img')
        for img in images:
            src = img.get('src')
            alt = img.get('alt', '')
            if src:
                if base_url:
                    src = urljoin(base_url, src)
                extracted_data['images'].append({
                    'src': src,
                    'alt': alt
                })
        
        # 5. Structured data (JSON-LD, microdata)
        json_ld_scripts = soup.find_all('script', type='application/ld+json')
        for script in json_ld_scripts:
            try:
                structured_data = json.loads(script.string)
                extracted_data['structured_data']['json_ld'] = structured_data
            except:
                pass
        
        # 6. Clean text extraction using html2text for better formatting
        h = html2text.HTML2Text()
        h.ignore_links = False
        h.ignore_images = False
        h.body_width = 0  # Don't wrap lines
        
        # Convert main content to clean text
        main_content_html = str(main_content)
        clean_text = h.handle(main_content_html)
        
        # Additional cleaning
        clean_text = re.sub(r'\n\s*\n\s*\n', '\n\n', clean_text)  # Remove excessive newlines
        clean_text = re.sub(r'[ \t]+', ' ', clean_text)  # Normalize spaces
        clean_text = clean_text.strip()
        
        extracted_data['clean_text'] = clean_text
        extracted_data['main_content'] = clean_text
        
        logger.info(f"HTML content extraction completed in {time.time() - start_time:.2f}s")
        logger.info(f"Extracted: {len(clean_text)} characters, {len(tables)} tables, {len(lists)} lists, {len(links)} links")
        
        return extracted_data
        
    except Exception as e:
        logger.error(f"HTML content extraction failed: {str(e)}")
        raise Exception(f"HTML content extraction failed: {str(e)}")

def format_html_content_for_processing(extracted_data: Dict[str, Any]) -> str:
    """
    Format extracted HTML data into a structured text format for document processing
    """
    formatted_parts = []
    
    # Add title
    if extracted_data.get('title'):
        formatted_parts.append(f"=== PAGE TITLE ===\n{extracted_data['title']}\n")
    
    # Add metadata if relevant
    metadata = extracted_data.get('metadata', {})
    relevant_meta = {}
    for key, value in metadata.items():
        if key.lower() in ['description', 'keywords', 'author', 'subject']:
            relevant_meta[key] = value
    
    if relevant_meta:
        formatted_parts.append("=== PAGE METADATA ===")
        for key, value in relevant_meta.items():
            formatted_parts.append(f"{key}: {value}")
        formatted_parts.append("")
    
    # Add structured data if available
    if extracted_data.get('structured_data'):
        formatted_parts.append("=== STRUCTURED DATA ===")
        structured_data = extracted_data['structured_data']
        formatted_parts.append(json.dumps(structured_data, indent=2))
        formatted_parts.append("")
    
    # Add tables with proper formatting
    tables = extracted_data.get('tables', [])
    for table in tables:
        formatted_parts.append(f"=== TABLE {table['table_number']} ===")
        
        if table['headers']:
            formatted_parts.append("HEADERS: " + " | ".join(table['headers']))
            formatted_parts.append("-" * 50)
        
        for row in table['rows']:
            formatted_parts.append(" | ".join(row))
        
        formatted_parts.append("")
    
    # Add lists
    lists = extracted_data.get('lists', [])
    for list_info in lists:
        formatted_parts.append(f"=== {list_info['type'].upper()} LIST {list_info['list_number']} ===")
        for i, item in enumerate(list_info['items'], 1):
            if list_info['type'] == 'ol':
                formatted_parts.append(f"{i}. {item}")
            else:
                formatted_parts.append(f"• {item}")
        formatted_parts.append("")
    
    # Add main content
    if extracted_data.get('main_content'):
        formatted_parts.append("=== MAIN CONTENT ===")
        formatted_parts.append(extracted_data['main_content'])
        formatted_parts.append("")
    
    # Add important links if they contain useful context
    links = extracted_data.get('links', [])
    important_links = [link for link in links if len(link['text']) > 10][:10]  # Limit to 10 most substantial links
    if important_links:
        formatted_parts.append("=== IMPORTANT LINKS ===")
        for link in important_links:
            formatted_parts.append(f"{link['text']}: {link['url']}")
        formatted_parts.append("")
    
    return "\n".join(formatted_parts)

# ---- Excel Parser with Proper Header-Value Association ----
def parse_excel_with_headers(file_content: bytes) -> str:
    """Parse Excel file with proper header-value association"""
    logger.info("Starting Excel parsing with header association")
    start_time = time.time()
    
    try:
        # Load workbook from bytes
        workbook = load_workbook(io.BytesIO(file_content), data_only=True)
        all_text_parts = []
        
        for sheet_name in workbook.sheetnames:
            logger.info(f"Processing Excel sheet: {sheet_name}")
            sheet = workbook[sheet_name]
            
            # Get sheet dimensions efficiently
            max_row = sheet.max_row
            max_col = sheet.max_column
            
            if max_row == 1 or max_col == 1:
                continue  # Skip empty or single-cell sheets
                
            # Convert to list of lists for easier processing
            data_rows = []
            for row in sheet.iter_rows(min_row=1, max_row=min(max_row, 1000), values_only=True):  # Limit rows for memory
                # Convert None to empty string and clean data
                cleaned_row = [str(cell).strip() if cell is not None else "" for cell in row]
                if any(cleaned_row):  # Only add non-empty rows
                    data_rows.append(cleaned_row)
            
            if len(data_rows) < 2:
                continue  # Need at least header + 1 data row
                
            # Identify header row (usually first row)
            headers = data_rows[0]
            data_rows = data_rows[1:]
            
            # Clean and validate headers
            clean_headers = []
            for i, header in enumerate(headers):
                if header and str(header).strip():
                    clean_headers.append(str(header).strip())
                else:
                    clean_headers.append(f"Column_{i+1}")  # Default name for empty headers
            
            # Build structured text with header-value pairs
            sheet_text_parts = [f"=== EXCEL SHEET: {sheet_name} ===\n"]
            
            # Method 1: Row-wise representation with clear header-value mapping
            for row_idx, row_data in enumerate(data_rows[:500], 1):  # Limit rows for memory
                if not any(str(cell).strip() for cell in row_data if cell):  # Skip empty rows
                    continue
                    
                row_text_parts = [f"Row {row_idx}:"]
                for header, value in zip(clean_headers, row_data):
                    if value and str(value).strip():
                        clean_value = str(value).strip()
                        # Create clear associations
                        row_text_parts.append(f"{header}: {clean_value}")
                
                if len(row_text_parts) > 1:  # Has actual data
                    sheet_text_parts.append(" | ".join(row_text_parts))
            
            # Method 2: Column-wise summary for better context
            sheet_text_parts.append(f"\n--- COLUMN SUMMARIES for {sheet_name} ---")
            for col_idx, header in enumerate(clean_headers):
                if col_idx >= len(data_rows[0]) if data_rows else True:
                    continue
                    
                # Extract column values
                column_values = []
                for row in data_rows[:100]:  # Sample first 100 rows
                    if col_idx < len(row) and row[col_idx] and str(row[col_idx]).strip():
                        val = str(row[col_idx]).strip()
                        if val not in column_values:  # Avoid duplicates
                            column_values.append(val)
                        if len(column_values) >= 10:  # Limit unique values shown
                            break
                
                if column_values:
                    values_preview = ", ".join(column_values[:5])
                    if len(column_values) > 5:
                        values_preview += f"... ({len(column_values)} unique values)"
                    sheet_text_parts.append(f"{header} contains: {values_preview}")
            
            all_text_parts.extend(sheet_text_parts)
            all_text_parts.append("\n" + "="*50 + "\n")
            
            # Memory cleanup with enhanced logging
            cleanup_variables(data_rows)
            clear_memory(f"Excel sheet {sheet_name}")
        
        workbook.close()
        final_text = "\n".join(all_text_parts)
        
        logger.info(f"Excel parsing completed in {time.time() - start_time:.2f}s. Text length: {len(final_text)}")
        return final_text
        
    except Exception as e:
        logger.error(f"Excel parsing failed: {str(e)}")
        raise Exception(f"Excel parsing failed: {str(e)}")

# ---- PowerPoint Parser ----
def parse_powerpoint(file_content: bytes) -> str:
    """Parse PowerPoint files extracting text and handling images with parallel OCR"""
    logger.info("Starting PowerPoint parsing with parallel image OCR support")
    start_time = time.time()
    
    try:
        presentation = Presentation(io.BytesIO(file_content))
        all_text_parts = []
        image_data_list = []  # Collect all images for parallel processing
        slide_image_mapping = {}  # Track which images belong to which slides
        
        # First pass: Extract text and collect images
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text_parts = [f"=== SLIDE {slide_num} ==="]
            slide_text = []
            
            for shape in slide.shapes:
                # Extract regular text
                if hasattr(shape, 'text') and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # Handle tables in slides
                if shape.has_table:
                    table = shape.table
                    table_text = ["TABLE DATA:"]
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_data.append(cell.text.strip())
                        if row_data:
                            table_text.append(" | ".join(row_data))
                    slide_text.append("\n".join(table_text))
                
                # Collect images for parallel OCR processing
                try:
                    if shape.shape_type == 13:  # Picture type
                        image_stream = shape.image.blob
                        image_name = f"slide_{slide_num}_image_{len(image_data_list)}"
                        image_data_list.append((image_stream, image_name))
                        slide_image_mapping[len(image_data_list) - 1] = slide_num
                        logger.info(f"Found image in slide {slide_num}, queued for parallel OCR")
                except Exception as img_error:
                    logger.warning(f"Could not extract image from slide {slide_num}: {str(img_error)}")
                    continue
            
            # Store slide text
            if slide_text:
                slide_text_parts.extend(slide_text)
            all_text_parts.append("\n".join(slide_text_parts))
        
        # Process all images in parallel
        if image_data_list:
            logger.info(f"Processing {len(image_data_list)} images with parallel OCR")
            ocr_results = process_images_parallel_ocr(image_data_list)
            
            # Add OCR results to respective slides
            for i, ocr_text in enumerate(ocr_results):
                if i < len(image_data_list):
                    slide_num = slide_image_mapping.get(i)
                    if slide_num and ocr_text:
                        # Find the slide in all_text_parts and append OCR text
                        for j, slide_content in enumerate(all_text_parts):
                            if f"=== SLIDE {slide_num} ===" in slide_content:
                                all_text_parts[j] += f"\nIMAGE_TEXT: {ocr_text}"
                                break
            
            logger.info(f"Parallel OCR completed, processed {len(ocr_results)} images successfully")
        
        final_text = "\n\n".join(all_text_parts)
        logger.info(f"PowerPoint parsing completed in {time.time() - start_time:.2f}s. Text length: {len(final_text)} characters")
        return final_text
        
    except Exception as e:
        logger.error(f"PowerPoint parsing failed: {str(e)}")
        raise Exception(f"PowerPoint parsing failed: {str(e)}")
    
def process_image_ocr_space(image_bytes: bytes, image_name: str) -> str:
    try:
        response = requests.post(
            url='https://api.ocr.space/parse/image',
            files={'filename': (image_name + '.jpg', image_bytes)},
            data={'apikey': OCR_SPACE_API_KEY, 'language': 'eng'}
        )
        result = response.json()
        return result['ParsedResults'][0]['ParsedText'] if 'ParsedResults' in result else ''
    except Exception as e:
        return f"[OCR Failed: {str(e)}]"

def process_images_parallel_ocr(image_data_list: List[Tuple[bytes, str]]) -> List[str]:
    """Process multiple images with OCR in parallel"""
    if not image_data_list:
        return []
    
    logger.info(f"Processing {len(image_data_list)} images with parallel OCR")
    ocr_results = []
    
    with ThreadPoolExecutor(max_workers=min(PARALLEL_OCR_BATCH, len(image_data_list))) as executor:
        # Submit all OCR tasks
        future_to_image = {
            executor.submit(process_image_ocr_space, img_bytes, img_name): (img_bytes, img_name)
            for img_bytes, img_name in image_data_list
        }
        
        # Collect results as they complete
        for future in as_completed(future_to_image):
            img_bytes, img_name = future_to_image[future]
            try:
                ocr_text = future.result()
                if ocr_text and ocr_text.strip() and "failed" not in ocr_text.lower():
                    ocr_results.append(ocr_text.strip())
                    logger.info(f"Successfully extracted text from {img_name}")
                else:
                    logger.info(f"No meaningful text found in {img_name}")
            except Exception as e:
                logger.warning(f"OCR failed for {img_name}: {str(e)}")
    
    return ocr_results
    
    
# ---- Word Document Parser ----
def parse_word_document(file_content: bytes) -> str:
    """Parse Word documents"""
    logger.info("Starting Word document parsing")
    start_time = time.time()
    
    try:
        document = Document(io.BytesIO(file_content))
        text_parts = []
        
        # Extract paragraphs
        for para in document.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())
        
        # Extract tables with proper structure
        for table in document.tables:
            table_text = ["TABLE DATA:"]
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_data.append(cell.text.strip())
                if row_data:
                    table_text.append(" | ".join(row_data))
            if len(table_text) > 1:
                text_parts.append("\n".join(table_text))
        
        final_text = "\n\n".join(text_parts)
        logger.info(f"Word document parsing completed in {time.time() - start_time:.2f}s. Text length: {len(final_text)}")
        return final_text
        
    except Exception as e:
        logger.error(f"Word document parsing failed: {str(e)}")
        raise Exception(f"Word document parsing failed: {str(e)}")

# ---- Image OCR Processing ----
def process_image_ocr(image_content: bytes, filename: str = "") -> str:
    """Process image with OCR, memory-optimized"""
    logger.info(f"Starting OCR processing for image: {filename}")
    start_time = time.time()
    
    try:
        # Load and optimize image
        image = Image.open(io.BytesIO(image_content))
        
        # Convert to RGB if needed
        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        # Resize if too large to save memory
        if image.size[0] > MAX_IMAGE_SIZE[0] or image.size[1] > MAX_IMAGE_SIZE[1]:
            image.thumbnail(MAX_IMAGE_SIZE, Image.Resampling.LANCZOS)
            logger.info(f"Resized image to {image.size} for memory efficiency")
        
        # Perform OCR with optimized settings
        ocr_config = '--oem 3 --psm 6 -l eng'  # Optimized OCR settings
        extracted_text = pytesseract.image_to_string(image, config=ocr_config)
        
        # Clean up image from memory
        image.close()
        cleanup_variables(image)
        clear_memory("OCR image processing")
        
        # Clean extracted text
        cleaned_text = re.sub(r'\s+', ' ', extracted_text.strip())
        
        logger.info(f"OCR completed in {time.time() - start_time:.2f}s. Text length: {len(cleaned_text)}")
        return cleaned_text if cleaned_text else "No text detected in image"
        
    except Exception as e:
        logger.error(f"OCR processing failed for {filename}: {str(e)}")
        return f"OCR processing failed for {filename}"

# ---- ZIP File Handler ----
async def extract_and_process_zip(file_content: bytes, depth: int = 0) -> str:
    """Extract and process ZIP files - go only one level deep, choosing first ZIP if multiple found"""
    MAX_ZIP_DEPTH = 1  # Only go one level deep
    
    if depth >= MAX_ZIP_DEPTH:
        logger.warning(f"ZIP nesting depth limit reached ({MAX_ZIP_DEPTH}). Skipping further extraction.")
        return "No meaningful content found - ZIP nesting too deep (max 1 level allowed)"
    
    logger.info(f"Starting ZIP file processing (depth level: {depth + 1})")
    start_time = time.time()
    
    all_extracted_text = []
    processed_files = []
    nested_zips_found = []
    first_nested_zip_processed = False
    
    try:
        with zipfile.ZipFile(io.BytesIO(file_content), 'r') as zip_file:
            file_list = zip_file.namelist()
            logger.info(f"Found {len(file_list)} files in ZIP")
            
            for file_name in file_list:
                # Skip directories and hidden files
                if file_name.endswith('/') or file_name.startswith('.'):
                    continue
                
                # Check if this is a ZIP file
                file_ext = get_file_extension(file_name)
                if file_ext == 'zip':
                    nested_zips_found.append(file_name)
                    # Only process the first nested ZIP found
                    if first_nested_zip_processed:
                        logger.info(f"Skipping nested ZIP {file_name} - already processed first nested ZIP")
                        all_extracted_text.append(f"=== FILE: {file_name} ===\nSkipped: Multiple nested ZIPs found, only processing first one\n")
                        continue
                
                # Check file size before extraction
                file_info = zip_file.getinfo(file_name)
                if file_info.file_size > MAX_FILE_SIZE:
                    logger.warning(f"Skipping {file_name}: file too large ({file_info.file_size} bytes)")
                    all_extracted_text.append(f"=== FILE: {file_name} ===\nSkipped: File too large ({file_info.file_size} bytes)\n")
                    continue
                
                # Check if format is supported
                is_supported, error_msg = is_supported_format(file_name)
                if not is_supported:
                    logger.info(f"Skipping {file_name}: {error_msg}")
                    all_extracted_text.append(f"=== FILE: {file_name} ===\nSkipped: {error_msg}\n")
                    continue
                
                try:
                    # Extract file content
                    extracted_content = zip_file.read(file_name)
                    logger.info(f"Processing extracted file: {file_name} (depth: {depth + 1})")
                    
                    # Special handling for nested ZIPs
                    if file_ext == 'zip':
                        if depth >= MAX_ZIP_DEPTH - 1:  # At maximum depth
                            logger.warning(f"Cannot process nested ZIP {file_name} - at maximum depth {depth + 1}")
                            all_extracted_text.append(f"=== FILE: {file_name} ===\nSkipped: ZIP nesting limit reached (max 1 level allowed)\n")
                            continue
                        else:
                            first_nested_zip_processed = True
                            logger.info(f"Processing first nested ZIP: {file_name}")
                    
                    # Process based on file type, passing depth for nested ZIPs
                    file_text = await process_file_content(extracted_content, file_name, depth)
                    
                    if file_text and file_text.strip():
                        all_extracted_text.append(f"=== FILE: {file_name} ===\n{file_text}\n")
                        processed_files.append(file_name)
                    else:
                        all_extracted_text.append(f"=== FILE: {file_name} ===\nNo content extracted\n")
                    
                    # Memory cleanup after each file
                    cleanup_variables(extracted_content)
                    clear_memory(f"ZIP file {file_name} processing")
                    
                except Exception as e:
                    logger.error(f"Error processing {file_name} from ZIP: {str(e)}")
                    all_extracted_text.append(f"=== FILE: {file_name} ===\nError processing file: {str(e)}\n")
        
        final_text = "\n".join(all_extracted_text)
        
        # Add summary information
        summary_info = []
        summary_info.append(f"=== ZIP PROCESSING SUMMARY ===")
        summary_info.append(f"Total files found: {len(file_list)}")
        summary_info.append(f"Successfully processed: {len(processed_files)}")
        
        if nested_zips_found:
            summary_info.append(f"Nested ZIPs found: {len(nested_zips_found)} ({', '.join(nested_zips_found[:5])}{' ...' if len(nested_zips_found) > 5 else ''})")
            if len(nested_zips_found) > 1:
                summary_info.append(f"Only processed first nested ZIP: {nested_zips_found[0]}")
                summary_info.append(f"Skipped nested ZIPs: {', '.join(nested_zips_found[1:])}")
        
        summary_info.append(f"Processing depth: {depth + 1}/{MAX_ZIP_DEPTH + 1}")
        summary_info.append(f"="*50)
        
        final_text = "\n".join(summary_info) + "\n\n" + final_text
        
        logger.info(f"📦 ZIP processing completed (depth {depth + 1}) in {time.time() - start_time:.2f}s. Processed {len(processed_files)} files")
        if nested_zips_found:
            logger.info(f"📦 Found {len(nested_zips_found)} nested ZIP(s), processed only first: {nested_zips_found[0] if nested_zips_found else 'none'}")
        
        return final_text if final_text.strip() else "No processable content found in ZIP file"
        
    except Exception as e:
        logger.error(f"ZIP processing failed: {str(e)}")
        raise Exception(f"ZIP processing failed: {str(e)}")

# ---- Unified File Content Processor ----
async def process_file_content(file_content: bytes, filename: str, depth: int = 0) -> str:
    """Process file content based on file type with depth tracking for nested ZIPs"""
    ext = get_file_extension(filename)
    
    try:
        if ext == 'pdf':
            # Use existing PDF processing
            doc = fitz.open(stream=file_content, filetype="pdf")
            text_parts = []
            max_pages = min(50, doc.page_count)
            
            for page_num in range(max_pages):
                page = doc[page_num]
                page_text = page.get_text()
                if page_text.strip():
                    page_text = words_to_numbers(page_text)
                    cleaned_text = re.sub(r'\s+', ' ', page_text.strip())
                    text_parts.append(cleaned_text)
            
            doc.close()
            return "\n\n".join(text_parts)
            
        elif ext in ['xlsx', 'xls']:
            return parse_excel_with_headers(file_content)
            
        elif ext in ['pptx', 'ppt']:
            return parse_powerpoint(file_content)
            
        elif ext in ['docx', 'doc']:
            return parse_word_document(file_content)
            
        elif ext in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp']:
            return process_image_ocr(file_content, filename)
            
        elif ext == 'txt':
            # Handle text files
            try:
                text_content = file_content.decode('utf-8')
            except UnicodeDecodeError:
                try:
                    text_content = file_content.decode('latin-1')
                except UnicodeDecodeError:
                    text_content = file_content.decode('utf-8', errors='ignore')
            
            return words_to_numbers(text_content.strip())
            
        elif ext == 'zip':
            return await extract_and_process_zip(file_content, depth + 1)
        
        
        elif ext in ['html', 'htm'] or 'html' in filename.lower():
            # Handle HTML files
            try:
                html_content = file_content.decode('utf-8')
            except UnicodeDecodeError:
                try:
                    html_content = file_content.decode('latin-1')
                except UnicodeDecodeError:
                    html_content = file_content.decode('utf-8', errors='ignore')
            
            # Extract and format HTML content
            extracted_data = extract_html_content(html_content)
            formatted_text = format_html_content_for_processing(extracted_data)
            return words_to_numbers(formatted_text.strip())
                
        else:
            raise Exception(f"Unsupported file format: {ext}")
            
    except Exception as e:
        logger.error(f"Error processing {filename}: {str(e)}")
        raise Exception(f"Error processing {filename}: {str(e)}")

# ---- Enhanced Document Downloader ----
async def download_and_process_document(url: str) -> str:
    """Download and process document from URL"""
    start_time = time.time()
    logger.info(f"Starting enhanced document download from: {url}")
    
    try:
        # Download with optimized settings
        timeout = httpx.Timeout(30.0, connect=5.0)
        limits = httpx.Limits(max_keepalive_connections=5, max_connections=10)
        
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        
        async with httpx.AsyncClient(timeout=timeout, limits=limits, follow_redirects=True, headers=headers) as client:
            response = await client.get(url)
        
        response.raise_for_status()
        
        # Check content type
        content_type = response.headers.get('content-type', '').lower()
        logger.info(f"Content type: {content_type}")
        
        # Check file size
        content_length = len(response.content)
        if content_length > MAX_FILE_SIZE:
            raise HTTPException(status_code=400, detail=f"File too large: {content_length} bytes. Maximum allowed: {MAX_FILE_SIZE} bytes")
        
        logger.info(f"Document downloaded successfully. Size: {content_length} bytes")
        
        # Determine processing method based on content type or URL
        if 'text/html' in content_type or url.endswith(('.html', '.htm')) or '<!DOCTYPE html' in response.text[:100].lower():
            # Process as HTML
            logger.info("Processing as HTML page")
            
            try:
                extracted_data = extract_html_content(response.text, url)
                document_text = format_html_content_for_processing(extracted_data)
                
                # Apply word-to-number conversion for consistency
                document_text = words_to_numbers(document_text)
                
                logger.info(f"HTML processing completed. Final text length: {len(document_text)}")
                return document_text
                
            except Exception as e:
                logger.error(f"HTML processing failed, trying as plain text: {str(e)}")
                # Fallback to plain text
                return words_to_numbers(response.text)
        
        else:
            # Process as file based on URL or content type
            filename = url.split('/')[-1].split('?')[0]  # Remove query parameters
            if not filename or '.' not in filename:
                # Try to determine from content type
                if 'pdf' in content_type:
                    filename = "document.pdf"
                elif 'excel' in content_type or 'spreadsheet' in content_type:
                    filename = "document.xlsx"
                elif 'powerpoint' in content_type or 'presentation' in content_type:
                    filename = "document.pptx"
                elif 'word' in content_type:
                    filename = "document.docx"
                else:
                    filename = "document.txt"
            
            # Check if format is supported
            is_supported, error_msg = is_supported_format(filename)
            if not is_supported and 'html' not in content_type:
                raise HTTPException(status_code=400, detail=error_msg)
            
            # Process based on file type
            document_text = await process_file_content(response.content, filename)
            
            logger.info(f"File processing completed. Text length: {len(document_text)}")
        
        # Clean up
        del response
        clear_memory()
        
        logger.info(f"Document processing completed in {time.time() - start_time:.2f}s. Text length: {len(document_text)}")
        return document_text
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Document download/processing failed: {str(e)}")
        raise Exception(f"Download or processing failed: {str(e)}")
# ---- Enhanced Embeddings and Search ----

async def get_embeddings_batch(client: httpx.AsyncClient, batch: List[str], input_type: str, batch_num: int) -> List[List[float]]:
    """Process a single batch of embeddings"""
    url = "https://api.cohere.com/v1/embed"
    current_cohere_key = get_current_cohere_key()
    headers = {
        "Authorization": f"Bearer {current_cohere_key}",
        "Content-Type": "application/json",
        "Accept": "application/json"
    }
    
    data = {
        "model": COHERE_MODEL,
        "texts": batch,
        "input_type": input_type,
        "truncate": "END"
    }
    
    try:
        response = await client.post(url, headers=headers, json=data)
        
        if response.status_code == 200:
            response_data = response.json()
            embeddings_data = response_data.get("embeddings", [])
            
            batch_embeddings = []
            for embedding in embeddings_data:
                vec = np.array(embedding, dtype=np.float32)
                norm = np.linalg.norm(vec)
                if norm > 0:
                    vec = vec / norm
                batch_embeddings.append(vec.tolist())
            
            logger.info(f"Completed embedding batch {batch_num}")
            return batch_embeddings
        else:
            logger.error(f"Cohere API error for batch {batch_num}: {response.status_code} - {response.text}")
            response.raise_for_status()
            
    except Exception as e:
        logger.error(f"Embedding batch {batch_num} failed: {str(e)}")
        raise e

async def get_embeddings(texts: List[str], input_type: str = "search_document") -> List[List[float]]:
    start_time = time.time()
    current_cohere_key = get_current_cohere_key()
    logger.info(f"🔑 Using Cohere API KEY #{current_cohere_key_index + 1} for {len(texts)} texts with parallel processing")
    
    clean_texts = []
    for text in texts:
        if text.strip():
            converted_text = words_to_numbers(text)
            clean_text = re.sub(r'\s+', ' ', converted_text.strip())[:EMBEDDING_MAX_LENGTH]
            clean_texts.append(clean_text)
    
    if not clean_texts:
        raise ValueError("No valid texts provided for embedding")
    
    # Create batches
    batches = [clean_texts[i:i+EMBEDDING_BATCH_SIZE] for i in range(0, len(clean_texts), EMBEDDING_BATCH_SIZE)]
    
    # Process batches with limited concurrency to avoid overwhelming the API
    timeout = httpx.Timeout(45.0)
    limits = httpx.Limits(max_keepalive_connections=5, max_connections=10)
    
    all_embeddings = []
    
    async with httpx.AsyncClient(timeout=timeout, limits=limits) as client:
        # Process batches in groups to control concurrency
        for i in range(0, len(batches), PARALLEL_EMBEDDING_CONCURRENT):
            concurrent_batches = batches[i:i+PARALLEL_EMBEDDING_CONCURRENT]
            
            # Create tasks for concurrent batch processing
            tasks = [
                get_embeddings_batch(client, batch, input_type, i + j + 1)
                for j, batch in enumerate(concurrent_batches)
            ]
            
            # Wait for all batches in this group to complete
            batch_results = await asyncio.gather(*tasks)
            
            # Collect results
            for batch_embeddings in batch_results:
                all_embeddings.extend(batch_embeddings)
    
    logger.info(f"Got {len(all_embeddings)} embeddings in {time.time() - start_time:.2f}s using parallel processing")
    return all_embeddings


def extract_key_terms(question: str) -> List[str]:
    important_terms = {
        'coverage', 'limit', 'deductible', 'premium', 'claim', 'benefit', 'exclusion',
        'copay', 'coinsurance', 'maximum', 'minimum', 'annual', 'lifetime', 'policy',
        'insured', 'covered', 'eligible', 'amount', 'percentage', 'network', 'provider',
        'emergency', 'prescription', 'medical', 'dental', 'vision', 'mental', 'health',
        'hospital', 'outpatient', 'inpatient', 'surgery', 'diagnostic', 'preventive',
        'waiting', 'period', 'authorization', 'existing', 'condition'
    }
    
    question_converted = words_to_numbers(question.lower())
    
    question_words = set(re.findall(r'\b\w+\b', question_converted))
    key_terms = list(important_terms.intersection(question_words))
    
    numbers = re.findall(r'\b\d+(?:\.\d+)?%?\b', question_converted)
    dollar_amounts = re.findall(r'\$[\d,]+(?:\.\d{2})?', question_converted)
    key_terms.extend(numbers + dollar_amounts)
    
    phrases = re.findall(r'\b(?:out of pocket|prior authorization|pre existing|waiting period|per year|per day)\b', question_converted)
    key_terms.extend(phrases)
    
    return key_terms

def is_yes_no_question(question: str) -> bool:
    question_lower = question.lower().strip()
    return question_lower.startswith(('is ', 'are ', 'do ', 'does ', 'did ', 'will ', 'would ', 'can ', 'could '))

def build_prompt(question: str, context_chunks: List[str]) -> str:
    context = "\n\n---\n\n".join(context_chunks[:15])  # Focused context for reasoning
    
    is_yes_no = is_yes_no_question(question)
    key_terms = extract_key_terms(question.lower())
    
    hints = ""
    if key_terms:
        hints = f"""
**KEY TERMS TO LOOK FOR:** {", ".join(key_terms[:15])}
(The above terms from your question should help identify relevant information in the context below)
"""
    # Language instruction - simple but effective

    lan = detect(question)  # Detect the language from the question
    lang = lan
    has_non_english = lan != "en"
    language_instruction = "**IMPORTANT: Respond in the SAME LANGUAGE as the question.**\n\n" if has_non_english else ""

    if is_yes_no:
        response_instruction = """This is a yes/no question. Start your response with "Yes" or "No", then provide a brief explanation with specific details in 25-40 words total."""
    else:
        response_instruction = """Answer in ONE concise, direct paragraph (40-80 words maximum). Include specific numbers, amounts, percentages, and conditions. Use figures (1, 2, 3) instead of words (one, two, three) for all numbers."""

    return f"""{language_instruction}You are an intelligent insurance analyst who understands policy documents and can reason about their content. Follow this answer hierarchy:

**ANSWER HIERARCHY (in order of priority):**
1. **DIRECT ANSWER**: If the exact answer is clearly stated in the context, provide it directly with specific references
2. **DEDUCED ANSWER**: If not directly stated, deduce the answer by connecting related information from the context
3. **INTELLIGENT REASONING**: If neither direct nor deduction is possible, provide a well-reasoned answer based on general insurance knowledge and principles

**YOUR TASK:** {response_instruction}

**REASONING APPROACH:**
1. First, search for DIRECT answers in the context - look for exact matches or explicit statements
2. If no direct answer, DEDUCE by finding patterns, rules, and logical connections in the context data
3. Look for mathematical patterns, sequences, or operations demonstrated in similar examples
4. Only use general knowledge if absolutely no patterns can be deduced from the context
5. Always be transparent about your reasoning: "According to the document...", "Following the pattern shown...", "Based on similar examples..."
6. Include specific numbers, dates, amounts, and conditions when available
7. PRIORITY: Pattern deduction from context over standard knowledge

**CONTEXT TO ANALYZE:**
{context}

**QUESTION:** {question}

**ANSWER (follow the hierarchy - direct → deduced → intelligent reasoning):**"""

async def ask_llm(prompt: str, retry_count: int = 0) -> str:
    try:
        current_api_key = get_current_gemini_key()
        logger.info(f"🔑 Using Gemini API KEY_{GEMINI_API_KEYS.index(current_api_key) + 1} (Request #{gemini_request_count}/12)")
        
        start_time = time.time()
        
        genai.configure(api_key=current_api_key)
        model = genai.GenerativeModel(GEMINI_MODEL)
        
        response = model.generate_content(
            prompt,
            generation_config=genai.types.GenerationConfig(
                temperature=LLM_TEMPERATURE,
                max_output_tokens=LLM_MAX_TOKENS,
                top_p=LLM_TOP_P,
                top_k=LLM_TOP_K,
            )
        )
        
        logger.info(f"Gemini API call completed in {time.time() - start_time:.2f}s")
        
        result_text = ""
        if hasattr(response, 'text') and response.text:
            result_text = response.text.strip()
        elif hasattr(response, 'candidates') and response.candidates:
            candidate = response.candidates[0]
            if hasattr(candidate, 'content') and candidate.content.parts:
                result_text = candidate.content.parts[0].text.strip()
        
        if result_text:
            result_text = words_to_numbers(result_text)
            result_text = re.sub(r'^(Answer:|ANSWER:|Response:|Based on the context:)\s*', '', result_text, flags=re.IGNORECASE)
            result_text = re.sub(r'\s+', ' ', result_text).strip()
            if result_text.endswith('.') and not result_text.endswith('etc.') and not result_text.endswith('vs.'):
                result_text = result_text[:-1]
            
            return result_text
        
        return "No valid response generated."
            
    except Exception as e:
        logger.error(f"Gemini API error (attempt {retry_count + 1}): {str(e)}")
        
        if retry_count < 1 and ("quota" not in str(e).lower() and "limit" not in str(e).lower()):
            logger.info("Retrying with next API key...")
            return await ask_llm(prompt, retry_count + 1)
        
        return f"Unable to generate answer due to API error. Please try again."
    
@app.get("/")
async def root():
    memory_info = get_memory_usage()
    return {
        "message": "Intelligent Multi-Format Document Processor with Parallel Processing & Reasoning", 
        "status": "intelligent_reasoning_v3.1",
        "supported_formats": list(SUPPORTED_FORMATS),
        "api_keys": {
            "gemini_keys_count": len(GEMINI_API_KEYS),
            "current_gemini_key": current_gemini_key_index + 1,
            "gemini_requests_on_current_key": gemini_request_count,
            "cohere_keys_count": len(COHERE_API_KEYS),
            "current_cohere_key": current_cohere_key_index + 1
        },
        "max_file_size_mb": MAX_FILE_SIZE / (1024 * 1024),
        "current_memory_usage_mb": memory_info.get('process_memory_mb', 0),
        "memory_optimization": "render_free_tier_ready",
        "endpoints": {
            "main": "/hackrx/run",
            "health": "/health", 
            "memory_status": "/memory-status",
            "api_status": "/api-status"
        }
    }

@app.post("/hackrx/run", response_model=QueryResponse)
async def run_query(request: QueryRequest, authorization: str = Header(None)):
    global qa_storage
    
    total_start_time = time.time()
    
    # Log initial memory state
    initial_memory = log_memory_usage("Request start")
    logger.info(f"📥 Received request with {len(request.questions)} questions")
    
    # Auth check
    if not authorization or not authorization.startswith("Bearer "):
        logger.error("Missing or malformed authorization header")
        raise HTTPException(status_code=401, detail="Authorization header missing or malformed.")

    token = authorization.split(" ")[1]
    if token != TEAM_TOKEN:
        logger.error("Invalid token provided")
        raise HTTPException(status_code=403, detail="Invalid token")

    try:
        # Process document with enhanced multi-format support
        logger.info("📄 Starting enhanced document processing")
        log_memory_usage("Before document processing")
        
        document_text = await download_and_process_document(request.documents)
        log_memory_usage("After document download")
        
        if not document_text.strip():
            logger.error("No text extracted from document")
            raise HTTPException(status_code=400, detail="No text could be extracted from the document")
        
        # Use improved semantic chunking
        logger.info("🧩 Starting document chunking")
        chunks = get_enhanced_chunks(document_text)
        log_memory_usage("After chunking")
        
        if not chunks:
            logger.error("No chunks created from document")
            raise HTTPException(status_code=400, detail="No meaningful chunks could be created from the document")
        
        # Clear document text from memory after chunking
        cleanup_variables(document_text)
        clear_memory("Document text cleanup")
        
        # Get embeddings for chunks and questions
        logger.info("🤖 Getting embeddings")
        chunk_texts = [chunk["text"] for chunk in chunks]
        chunk_embeddings = await get_embeddings(chunk_texts, input_type="search_document")
        log_memory_usage("After chunk embeddings")
        
        question_embeddings = await get_embeddings(request.questions, input_type="search_query")
        log_memory_usage("After question embeddings")
        
        # Cleanup intermediate variables
        cleanup_variables(chunk_texts)
        clear_memory("Embeddings cleanup")
        
        # Create FAISS index with chunk embeddings
        logger.info("📊 Creating FAISS index")
        faiss_service.create_index(chunk_embeddings, chunks)
        log_memory_usage("After FAISS index creation")
        
        # Process questions
        logger.info("❓ Processing questions")
        answers = []
        for i, (question, q_emb) in enumerate(zip(request.questions, question_embeddings)):
            try:
                logger.info(f"📝 Processing question {i+1}/{len(request.questions)}: {question[:50]}...")
                question_start_memory = log_memory_usage(f"Question {i+1} start")
                
                # Use multi-tier FAISS search for comprehensive results
                search_results = faiss_service.multi_tier_search(q_emb, question, FAISS_K_SEARCH)
                
                # Optional enhanced search (only if USE_ENHANCED_SEARCH=true)
                from app.services.enhanced_search import apply_enhanced_search_if_enabled
                search_results = apply_enhanced_search_if_enabled(search_results, question)
                
                if not search_results:
                    logger.warning(f"No search results found for question: {question}")
                    # Use intelligent reasoning even without search results
                    fallback_prompt = f"""You are an intelligent analyst. No specific context was found, but try to reason logically about the question.

**QUESTION:** {question}

**INSTRUCTIONS:**
- Look for mathematical patterns, logical sequences, or domain-specific rules
- If it's a mathematical question, try to find underlying logic or patterns
- Only use general knowledge if no logical patterns can be determined
- Be transparent about your reasoning approach
- Keep answer concise (40-80 words)
- Use figures (1, 2, 3) instead of words for numbers

**ANSWER (logical reasoning or general knowledge):**"""
                    
                    try:
                        answer = await ask_llm(fallback_prompt)
                        if answer and not answer.startswith("Unable to generate"):
                            # Prefix to indicate this is based on general knowledge
                            answer = f"Based on general insurance principles: {answer}"
                        else:
                            answer = "I'd need more specific policy information to provide a definitive answer for this question."
                    except Exception as e:
                        logger.error(f"Fallback answer generation failed: {str(e)}")
                        answer = "I'd need more specific policy information to provide a definitive answer for this question."
                    
                    answers.append(answer)
                    qa_storage.append([question, answer])
                    continue
                
                # Check if search results have good relevance scores
                best_score = search_results[0].get('final_score', 0) if hasattr(search_results[0], 'get') and 'final_score' in search_results[0] else search_results[0].get('similarity_score', 0)
                
                if best_score < 0.15:  # Low relevance threshold
                    logger.info(f"Low relevance score ({best_score:.3f}) for question: {question[:50]}... Using hybrid approach")
                    
                    # Combine context with intelligent reasoning
                    relevant_chunks = faiss_service.enhance_results(search_results[:5], question)  # Use fewer chunks
                    limited_context = "\n\n---\n\n".join(relevant_chunks[:5])
                    
                    hybrid_prompt = f"""You are an intelligent data analyst. The available context has limited relevance, but analyze it carefully for patterns, rules, or examples.

**APPROACH:**
1. Look for ANY patterns, sequences, or mathematical operations in the context
2. Find similar examples and deduce the underlying rule or logic
3. Apply the discovered pattern to answer the question
4. If no patterns exist, then provide reasoning based on the domain
5. Be transparent: "Following the pattern...", "Based on similar examples...", "The data shows..."

**CONTEXT TO ANALYZE FOR PATTERNS:**
{limited_context}

**QUESTION:** {question}

**ANSWER (look for patterns first, then reasoning):**"""
                    
                    response = await ask_llm(hybrid_prompt)
                else:
                    # Good relevance - use normal approach
                    relevant_chunks = faiss_service.enhance_results(search_results, question)
                    prompt = build_prompt(question, relevant_chunks)
                    response = await ask_llm(prompt)
                
                # Clean response
                response = response.strip()
                response = re.sub(r'^(Answer:|ANSWER:|Response:)\s*', '', response, flags=re.IGNORECASE)
                response = re.sub(r'\s*\.$', '', response)
                
                final_answer = response.strip()
                answers.append(final_answer)
                
                # Store Q&A pair
                qa_storage.append([question, final_answer])
                
                # Cleanup question-specific variables
                cleanup_variables(search_results, relevant_chunks, prompt, response)
                
                logger.info(f"✅ Question {i+1} processed successfully")
                clear_memory(f"Question {i+1} completion")
                log_memory_usage(f"Question {i+1} end")
                
            except Exception as e:
                logger.error(f"Error processing question {i+1} '{question}': {str(e)}")
                error_answer = "Error processing the question. Please try again."
                answers.append(error_answer)
                qa_storage.append([question, error_answer])
        
        # Final cleanup of major variables
        cleanup_variables(chunks, chunk_embeddings, question_embeddings)
        
        total_time = time.time() - total_start_time
        final_memory = log_memory_usage("Request completion")
        
        # Calculate memory usage statistics
        if initial_memory and final_memory:
            memory_delta = final_memory['process_memory_mb'] - initial_memory['process_memory_mb']
            logger.info(f"📊 Memory Delta: {memory_delta:+.2f}MB from start to finish")
        
        logger.info(f"⏱️ Total request processed in {total_time:.2f}s")
        
        # Log Q&A pairs and clear storage
        logger.info("📋 ALL QUESTIONS AND ANSWERS:")
        logger.info(f"{qa_storage}")
        qa_storage.clear()
        
        # Final comprehensive memory cleanup
        clear_memory("Complete request cleanup")
        
        # Rotate Cohere API key for the next request
        rotate_cohere_key()
        
        return QueryResponse(answers=answers)
        
    except HTTPException as he:
        logger.error(f"HTTP Exception: {he.detail}")
        raise he
    except Exception as e:
        logger.error(f"Unexpected error in run_query: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Internal server error: {str(e)}")

@app.get("/health")
async def health_check():
    faiss_stats = faiss_service.get_stats()
    memory_info = get_memory_usage()
    return {
        "status": "healthy", 
        "memory_usage": memory_info, 
        "version": "enhanced_memory_optimized_v2.1",
        "supported_formats": list(SUPPORTED_FORMATS),
        "unsupported_formats": list(UNSUPPORTED_FORMATS),
        "gemini_keys_available": len(GEMINI_API_KEYS),
        "current_gemini_key": current_gemini_key_index + 1,
        "gemini_requests_on_current_key": gemini_request_count,
        "cohere_keys_available": len(COHERE_API_KEYS),
        "current_cohere_key": current_cohere_key_index + 1,
        "qa_storage_size": len(qa_storage),
        "max_file_size_mb": MAX_FILE_SIZE / (1024 * 1024),
        "faiss_index": faiss_stats,
        "enhanced_features": {
            "semantic_chunking": True,
            "topic_based_chunking": True,
            "hybrid_search": True,
            "faiss_indexing": True,
            "quality_scoring": True,
            "parallel_chunking": True,
            "parallel_embedding": True,
            "parallel_ocr": True,
            "intelligent_reasoning": True,
            "context_deduction": True,
            "fallback_intelligence": True
        },
        "parallel_processing": {
            "max_workers": MAX_WORKERS,
            "parallel_chunk_size": PARALLEL_CHUNK_SIZE,
            "parallel_ocr_batch": PARALLEL_OCR_BATCH,
            "parallel_embedding_concurrent": PARALLEL_EMBEDDING_CONCURRENT
        }
    }

@app.get("/memory-status")
async def memory_status():
    """Endpoint to check current memory usage and perform cleanup"""
    # Get memory before cleanup
    memory_before = get_memory_usage()
    
    # Force cleanup
    collected = clear_memory("Memory status check")
    
    # Get memory after cleanup
    memory_after = get_memory_usage()
    
    return {
        "memory_before_cleanup": memory_before,
        "memory_after_cleanup": memory_after,
        "objects_collected": collected,
        "cleanup_performed": True,
        "memory_freed_mb": round(memory_before.get('process_memory_mb', 0) - memory_after.get('process_memory_mb', 0), 2) if memory_before and memory_after else 0,
        "render_free_tier_optimized": True
    }

@app.get("/api-status")
async def api_status():
    """Endpoint to check API key rotation status"""
    return {
        "total_api_keys": len(GEMINI_API_KEYS),
        "current_key_index": current_key_index + 1,
        "requests_on_current_key": request_count,
        "requests_per_key_limit": REQUESTS_PER_KEY,
        "qa_storage_current_size": len(qa_storage),
        "supported_formats": list(SUPPORTED_FORMATS),
        "memory_limit_mb": MAX_FILE_SIZE / (1024 * 1024)
    }

@app.get("/faiss-stats")
async def faiss_statistics():
    """Endpoint to check FAISS index statistics"""
    stats = faiss_service.get_stats()
    return {
        "faiss_index_info": stats,
        "search_capabilities": {
            "semantic_search": True,
            "keyword_search": stats.get("tfidf_available", False),
            "hybrid_search": True,
            "domain_specific_boosting": True,
            "numeric_matching": True,
            "phrase_matching": True
        },
        "chunking_strategies": {
            "sliding_window": True,
            "semantic_chunking": True,
            "topic_based_chunking": True,
            "structured_chunking": True,
            "quality_ranking": True
        }
    }

@app.get("/supported-formats")
async def supported_formats():
    """Endpoint to check supported file formats"""
    return {
        "supported_formats": {
            "documents": ["pdf", "docx", "doc", "txt"],
            "spreadsheets": ["xlsx", "xls"],
            "presentations": ["pptx", "ppt"],
            "images": ["jpg", "jpeg", "png", "gif", "bmp", "tiff", "webp"],
            "archives": ["zip"]
        },
        "unsupported_formats": list(UNSUPPORTED_FORMATS),
        "max_file_size_mb": MAX_FILE_SIZE / (1024 * 1024),
        "features": {
            "excel_header_mapping": True,
            "nested_zip_extraction": True,
            "image_ocr": True,
            "memory_optimized": True,
            "format_validation": True,
            "faiss_search": True,
            "hybrid_retrieval": True,
            "semantic_chunking": True,
            "topic_based_chunking": True,
            "quality_scoring": True,
            "enhanced_precision": True
        }
    }

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000, log_level="info")