# app/processors/powerpoint_processor.py
import io
import time
import logging
from typing import List, Tuple
from pptx import Presentation
from concurrent.futures import ThreadPoolExecutor, as_completed
from app.processors.base_processor import BaseProcessor
from app.config.settings import PARALLEL_OCR_BATCH
from app.processors.image_processor import process_image_ocr_space

logger = logging.getLogger(__name__)

class PowerPointProcessor(BaseProcessor):
    """PowerPoint document processor with parallel OCR support"""
    
    def get_supported_extensions(self) -> set:
        return {'pptx', 'ppt'}
    
    def supports_format(self, file_format: str) -> bool:
        """Check if the processor supports the given file format"""
        return file_format.lower() in self.get_supported_extensions()
    
    async def process(self, file_content: bytes, filename: str, **kwargs) -> str:
        """Parse PowerPoint files extracting text and handling images with parallel OCR"""
        logger.info("Starting PowerPoint parsing with parallel image OCR support")
        start_time = time.time()
        
        try:
            presentation = Presentation(io.BytesIO(file_content))
            all_text_parts = []
            image_data_list = []  # Collect all images for parallel processing
            slide_image_mapping = {}  # Track which images belong to which slides
            
            # First pass: Extract text and collect images
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text_parts = [f"=== SLIDE {slide_num} ==="]
                slide_text = []
                
                for shape in slide.shapes:
                    # Extract regular text
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Handle tables in slides
                    if shape.has_table:
                        table = shape.table
                        table_text = ["TABLE DATA:"]
                        for row in table.rows:
                            row_data = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_data.append(cell.text.strip())
                            if row_data:
                                table_text.append(" | ".join(row_data))
                        slide_text.append("\n".join(table_text))
                    
                    # Collect images for parallel OCR processing
                    try:
                        if shape.shape_type == 13:  # Picture type
                            image_stream = shape.image.blob
                            image_name = f"slide_{slide_num}_image_{len(image_data_list)}"
                            image_data_list.append((image_stream, image_name))
                            slide_image_mapping[len(image_data_list) - 1] = slide_num
                            logger.info(f"Found image in slide {slide_num}, queued for parallel OCR")
                    except Exception as img_error:
                        logger.warning(f"Could not extract image from slide {slide_num}: {str(img_error)}")
                        continue
                
                # Store slide text
                if slide_text:
                    slide_text_parts.extend(slide_text)
                all_text_parts.append("\n".join(slide_text_parts))
            
            # Process all images in parallel
            if image_data_list:
                logger.info(f"Processing {len(image_data_list)} images with parallel OCR")
                ocr_results = self._process_images_parallel_ocr(image_data_list)
                
                # Add OCR results to respective slides
                for i, ocr_text in enumerate(ocr_results):
                    if i < len(image_data_list):
                        slide_num = slide_image_mapping.get(i)
                        if slide_num and ocr_text:
                            # Find the slide in all_text_parts and append OCR text
                            for j, slide_content in enumerate(all_text_parts):
                                if f"=== SLIDE {slide_num} ===" in slide_content:
                                    all_text_parts[j] += f"\nIMAGE_TEXT: {ocr_text}"
                                    break
                
                logger.info(f"Parallel OCR completed, processed {len(ocr_results)} images successfully")
            
            final_text = "\n\n".join(all_text_parts)
            logger.info(f"PowerPoint parsing completed in {time.time() - start_time:.2f}s. Text length: {len(final_text)} characters")
            return final_text
            
        except Exception as e:
            logger.error(f"PowerPoint parsing failed: {str(e)}")
            raise Exception(f"PowerPoint parsing failed: {str(e)}")
    
    def _process_images_parallel_ocr(self, image_data_list: List[Tuple[bytes, str]]) -> List[str]:
        """Process multiple images with OCR in parallel"""
        if not image_data_list:
            return []
        
        logger.info(f"Processing {len(image_data_list)} images with parallel OCR")
        ocr_results = []
        
        with ThreadPoolExecutor(max_workers=min(PARALLEL_OCR_BATCH, len(image_data_list))) as executor:
            # Submit all OCR tasks
            future_to_image = {
                executor.submit(process_image_ocr_space, img_bytes, img_name): (img_bytes, img_name)
                for img_bytes, img_name in image_data_list
            }
            
            # Collect results as they complete
            for future in as_completed(future_to_image):
                img_bytes, img_name = future_to_image[future]
                try:
                    ocr_text = future.result()
                    if ocr_text and ocr_text.strip() and "failed" not in ocr_text.lower():
                        ocr_results.append(ocr_text.strip())
                        logger.info(f"Successfully extracted text from {img_name}")
                    else:
                        logger.info(f"No meaningful text found in {img_name}")
                except Exception as e:
                    logger.warning(f"OCR failed for {img_name}: {str(e)}")
        
        return ocr_results